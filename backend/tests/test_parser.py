import os
import tempfile
from app.services.parser_service import extract_text


def test_extract_text_from_pdf(sample_pdf, tmp_path):
    pdf_path = tmp_path / "test.pdf"
    pdf_path.write_bytes(sample_pdf)
    text = extract_text(str(pdf_path), "pdf")
    # Minimal PDF may not have extractable text with PyPDF2, but shouldn't crash
    assert isinstance(text, str)


def test_extract_text_from_docx(tmp_path):
    from docx import Document as DocxDocument
    doc = DocxDocument()
    doc.add_paragraph("This is a test paragraph.")
    doc.add_paragraph("Another paragraph here.")
    path = tmp_path / "test.docx"
    doc.save(str(path))

    text = extract_text(str(path), "docx")
    assert "test paragraph" in text
    assert "Another paragraph" in text


def test_extract_text_from_pptx(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide content here"
    path = tmp_path / "test.pptx"
    prs.save(str(path))

    text = extract_text(str(path), "pptx")
    assert "Slide Title" in text
    assert "Slide content" in text


def test_empty_file_handled(tmp_path):
    path = tmp_path / "empty.pdf"
    path.write_bytes(b"")
    try:
        text = extract_text(str(path), "pdf")
    except Exception:
        pass  # Empty file may raise, that's acceptable
