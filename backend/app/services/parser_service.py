import os
import subprocess
import tempfile
import glob as globmod
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


def extract_text(file_path: str, file_type: str) -> str:
    dispatch = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
        "doc": _extract_doc,
        "ppt": _extract_ppt,
        "xlsx": _extract_xlsx,
        "xls": _extract_xls,
        "txt": _extract_txt,
        "md": _extract_txt,
    }
    extractor = dispatch.get(file_type)
    if extractor:
        try:
            return extractor(file_path)
        except Exception:
            return ""
    return ""


def _extract_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n".join(texts)


def _extract_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            texts.append(para.text)
    return "\n".join(texts)


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return "\n".join(texts)


def _extract_doc(file_path: str) -> str:
    try:
        doc = DocxDocument(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception:
        pass
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "docx", "--outdir", tmpdir, file_path],
                capture_output=True, timeout=60,
            )
            docx_files = globmod.glob(os.path.join(tmpdir, "*.docx"))
            if docx_files:
                return _extract_docx(docx_files[0])
    except Exception:
        pass
    return ""


def _extract_ppt(file_path: str) -> str:
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pptx", "--outdir", tmpdir, file_path],
                capture_output=True, timeout=60,
            )
            pptx_files = globmod.glob(os.path.join(tmpdir, "*.pptx"))
            if pptx_files:
                return _extract_pptx(pptx_files[0])
    except Exception:
        pass
    return ""


def _extract_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    texts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                texts.append(row_text)
    return "\n".join(texts)


def _extract_xls(file_path: str) -> str:
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "xlsx", "--outdir", tmpdir, file_path],
                capture_output=True, timeout=60,
            )
            xlsx_files = globmod.glob(os.path.join(tmpdir, "*.xlsx"))
            if xlsx_files:
                return _extract_xlsx(xlsx_files[0])
    except Exception:
        pass
    return ""


def _extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
